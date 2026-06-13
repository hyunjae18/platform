"""
converter.py
────────────
Converts any uploaded document into a list of PIL Images,
one image per page. EasyOCR / PaddleOCR then processes each image.

Supported formats:
  Images : JPEG, PNG, TIFF, BMP, WebP, GIF
  PDF    : via PyMuPDF (fitz) — renders each page at configurable DPI (default 300)
  DOCX   : via python-docx — extracts embedded images + renders text as image if none
  PPTX   : via python-pptx — renders each slide as image
"""

import io
import logging
from pathlib import Path
from typing import Optional, List, Tuple

from PIL import Image

logger = logging.getLogger(__name__)

# ── Supported MIME types grouped by handler ───────────────────────────────────

IMAGE_MIMES = {
    "image/jpeg", "image/jpg", "image/png", "image/tiff",
    "image/bmp", "image/webp", "image/gif",
}

PDF_MIMES = {
    "application/pdf",
}

DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}

PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}

ALL_SUPPORTED_MIMES = IMAGE_MIMES | PDF_MIMES | DOCX_MIMES | PPTX_MIMES

EXTENSION_MIME_MAP = {
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png":  "image/png",
    ".tiff": "image/tiff",
    ".tif":  "image/tiff",
    ".bmp":  "image/bmp",
    ".webp": "image/webp",
    ".gif":  "image/gif",
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc":  "application/msword",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt":  "application/vnd.ms-powerpoint",
}


def resolve_mime(filename: str, declared_mime: Optional[str]) -> str:
    """Return the best MIME type for a file, falling back to extension."""
    if declared_mime and declared_mime in ALL_SUPPORTED_MIMES:
        return declared_mime
    ext = Path(filename).suffix.lower()
    return EXTENSION_MIME_MAP.get(ext, declared_mime or "application/octet-stream")


# ── Converters (return PIL Images) ───────────────────────────────────────────

def _image_bytes_to_pil(data: bytes) -> List[Image.Image]:
    """Plain image file → single-element list."""
    image = Image.open(io.BytesIO(data)).convert("RGB")
    return [image]


def _pdf_to_images(data: bytes, dpi: int = 300) -> List[Image.Image]:
    """PDF → one PIL Image per page using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("PyMuPDF required for PDF support. Run: pip install pymupdf")

    images = []
    doc = fitz.open(stream=data, filetype="pdf")
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        images.append(img)
        logger.info(f"PDF page {page_num+1}/{len(doc)} → {img.size}")

    doc.close()
    logger.info(f"PDF converted: {len(images)} pages at {dpi} DPI")
    return images


def _docx_to_images(data: bytes, dpi: int = 200) -> List[Image.Image]:
    """DOCX → embedded images; if none, render text as image."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx required. Run: pip install python-docx")

    images = []
    doc = Document(io.BytesIO(data))

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data)).convert("RGB")
                images.append(img)
                logger.info(f"DOCX embedded image extracted: {img.size}")
            except Exception as e:
                logger.warning(f"Could not extract DOCX image: {e}")

    if not images:
        logger.info("No embedded images in DOCX — rendering text as image.")
        images = _render_text_as_image(doc)

    return images


def _render_text_as_image(doc, chars_per_line: int = 80, font_size: int = 20) -> List[Image.Image]:
    """Render DOCX paragraphs onto white PIL Images."""
    try:
        from PIL import ImageDraw, ImageFont
    except ImportError:
        raise RuntimeError("Pillow required.")

    all_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if not all_text.strip():
        return []

    words = all_text.split()
    lines = []
    cur_line = []
    for word in words:
        cur_line.append(word)
        if len(" ".join(cur_line)) >= chars_per_line:
            lines.append(" ".join(cur_line))
            cur_line = []
    if cur_line:
        lines.append(" ".join(cur_line))

    lines_per_page = 50
    pages = [lines[i:i+lines_per_page] for i in range(0, max(len(lines), 1), lines_per_page)]

    images = []
    W, H = 1200, 1600
    for page_lines in pages:
        img = Image.new("RGB", (W, H), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        y = 40
        for line in page_lines:
            draw.text((40, y), line, fill=(0, 0, 0))
            y += font_size + 8
            if y > H - 40:
                break
        images.append(img)
    return images


def _pptx_to_images(data: bytes, dpi: int = 150) -> List[Image.Image]:
    """PPTX → one image per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx required. Run: pip install python-pptx")

    images = []
    prs = Presentation(io.BytesIO(data))

    for slide_num, slide in enumerate(prs.slides):
        slide_images = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                try:
                    img_data = shape.image.blob
                    img = Image.open(io.BytesIO(img_data)).convert("RGB")
                    slide_images.append(img)
                except Exception as e:
                    logger.warning(f"PPTX slide {slide_num+1} image error: {e}")

        if slide_images:
            images.extend(slide_images)
        else:
            # Text-only slide → render as image
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                from PIL import ImageDraw
                img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
                draw = ImageDraw.Draw(img)
                y = 40
                for text in texts:
                    words = text.split()
                    line = []
                    for word in words:
                        line.append(word)
                        if len(" ".join(line)) > 60:
                            draw.text((40, y), " ".join(line), fill=(0, 0, 0))
                            y += 28
                            line = []
                    if line:
                        draw.text((40, y), " ".join(line), fill=(0, 0, 0))
                        y += 28
                    y += 10
                images.append(img)

    logger.info(f"PPTX converted: {len(images)} images from {len(prs.slides)} slides")
    return images


# ── Public API (kept as original function name) ─────────────────────────────
def convert_to_images(
    file_bytes: bytes,
    filename: str,
    mime_type: Optional[str] = None,
    dpi: int = 300,
) -> Tuple[List[Image.Image], str]:
    """
    Convert any supported document to a list of PIL Images.

    Returns:
        (images, resolved_mime)
        images        : list of PIL Images, one per page/slide
        resolved_mime : the MIME type actually used

    Raises:
        ValueError  : unsupported file type
        RuntimeError: missing dependency
    """
    resolved = resolve_mime(filename, mime_type)
    logger.info(f"Converting '{filename}' (mime={resolved}, {len(file_bytes):,} bytes)")

    if resolved in IMAGE_MIMES:
        return _image_bytes_to_pil(file_bytes), resolved

    if resolved in PDF_MIMES:
        return _pdf_to_images(file_bytes, dpi=dpi), resolved

    if resolved in DOCX_MIMES:
        return _docx_to_images(file_bytes, dpi=dpi), resolved

    if resolved in PPTX_MIMES:
        return _pptx_to_images(file_bytes, dpi=dpi), resolved

    raise ValueError(
        f"Unsupported file type '{resolved}' for file '{filename}'. "
        f"Supported: images (JPEG/PNG/TIFF/BMP/WebP), PDF, DOCX, PPTX"
    )